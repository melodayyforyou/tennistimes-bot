"""
Tennistimes.id Python Worker
Flask app with 3 file-generation endpoints:

  POST /generate-pptx  — JSON → PowerPoint (.pptx)
    Brand: black background, gold titles, white body text, TennisTV.id footer

  POST /generate-xlsx  — JSON → Excel (.xlsx)
    Brand: green headers (#1D6A3A), white header text, TennisTV.id in A1

  POST /generate-pdf   — JSON → PDF report
    Brand: TennisTV.id header, bold section headings, page numbers in footer

Deployed on Render. Called by the Node.js bot via axios.
"""

import io
import os

from flask import Flask, request, send_file, jsonify
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment
from fpdf import FPDF

app = Flask(__name__)


# ══════════════════════════════════════════════════════════════════════════════
#  PPTX Generator
#  Brand spec: black background, gold/amber title (#F5A623), white body,
#              TennisTV.id footer on every slide
# ══════════════════════════════════════════════════════════════════════════════

@app.route('/generate-pptx', methods=['POST'])
def generate_pptx():
    """
    Expected JSON (matches system prompt schema):
    {
      "title": "deck title",
      "slides": [
        { "title": "slide title", "content": ["bullet 1", "bullet 2"], "notes": "speaker notes" }
      ]
    }
    """
    data = request.get_json(silent=True)
    if not data:
        return jsonify({'error': 'No valid JSON in request body'}), 400

    try:
        prs = Presentation()
        prs.slide_width  = Inches(13.33)  # 16:9 widescreen
        prs.slide_height = Inches(7.5)

        # Brand colours
        BLACK  = RGBColor(0x00, 0x00, 0x00)
        GOLD   = RGBColor(0xF5, 0xA6, 0x23)
        WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
        GRAY   = RGBColor(0xAA, 0xAA, 0xAA)

        def safe_text(text):
            """Strip characters that pptx XML can't encode."""
            import re
            text = str(text or '')
            # Remove control characters except tab/newline
            return re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)

        slides_data = data.get('slides', [])

        for i, slide_info in enumerate(slides_data):
            is_title_slide = (i == 0)

            # Layout 0 = Title Slide, Layout 1 = Title + Content
            layout = prs.slide_layouts[0] if is_title_slide else prs.slide_layouts[1]
            slide  = prs.slides.add_slide(layout)

            # ── Black background ────────────────────────────────────────────
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = BLACK

            # ── Slide title in gold ─────────────────────────────────────────
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = safe_text(slide_info.get('title', ''))
                for para in title_shape.text_frame.paragraphs:
                    para.font.bold  = True
                    para.font.color.rgb = GOLD
                    para.font.size  = Pt(40) if is_title_slide else Pt(30)

            # ── Deck title as subtitle on slide 1 ──────────────────────────
            if is_title_slide:
                for ph in slide.placeholders:
                    if ph.placeholder_format.idx == 1:
                        ph.text = safe_text(data.get('title', ''))
                        for para in ph.text_frame.paragraphs:
                            para.font.size      = Pt(22)
                            para.font.color.rgb = GRAY
                        break

            # ── Bullet content on body slides ───────────────────────────────
            else:
                content = slide_info.get('content', [])
                if content:
                    for ph in slide.placeholders:
                        if ph.placeholder_format.idx == 1:
                            tf = ph.text_frame
                            tf.clear()
                            for j, bullet in enumerate(content):
                                para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                                para.text  = safe_text(bullet)
                                para.level = 0
                                para.font.size      = Pt(20)
                                para.font.color.rgb = WHITE
                            break

            # ── Speaker notes ───────────────────────────────────────────────
            notes_text = slide_info.get('notes', '')
            if notes_text:
                slide.notes_slide.notes_text_frame.text = notes_text

            # ── TennisTV.id footer text box ─────────────────────────────────
            footer = slide.shapes.add_textbox(
                Inches(0.3),
                Inches(7.1),
                Inches(12.7),
                Inches(0.3),
            )
            tf = footer.text_frame
            tf.text = 'TennisTV.id'
            para    = tf.paragraphs[0]
            para.font.size      = Pt(9)
            para.font.color.rgb = GRAY
            para.font.italic    = True
            para.alignment      = PP_ALIGN.RIGHT

        output = io.BytesIO()
        prs.save(output)
        output.seek(0)

        return send_file(
            output,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name='presentation.pptx',
        )

    except Exception as exc:
        app.logger.error(f'/generate-pptx error: {exc}')
        return jsonify({'error': str(exc)}), 500


# ══════════════════════════════════════════════════════════════════════════════
#  XLSX Generator
#  Brand spec: TennisTV.id in A1, green header fill (#1D6A3A), white header
#              text, auto-sized columns, frozen header row
# ══════════════════════════════════════════════════════════════════════════════

@app.route('/generate-xlsx', methods=['POST'])
def generate_xlsx():
    """
    Expected JSON (matches system prompt schema — flat, single sheet):
    {
      "title": "sheet title",
      "headers": ["col1", "col2", "col3"],
      "rows": [["val1", "val2", "val3"]],
      "summary": "one line description"
    }
    """
    data = request.get_json(silent=True)
    if not data:
        return jsonify({'error': 'No valid JSON in request body'}), 400

    try:
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = 'Data'

        GREEN_FILL  = PatternFill(start_color='1D6A3A', end_color='1D6A3A', fill_type='solid')
        WHITE_FONT  = Font(bold=True, color='FFFFFF', size=12)
        TITLE_FONT  = Font(bold=True, size=14, color='1D6A3A')
        CENTER      = Alignment(horizontal='center', vertical='center')

        title   = data.get('title', 'Sheet')
        headers = data.get('headers', [])
        rows    = data.get('rows', [])

        def safe_cell(v):
            """Return numeric types as-is; coerce everything else to a plain string."""
            if isinstance(v, (int, float)):
                return v
            return str(v) if v is not None else ''

        # Row 1 — TennisTV.id title spanning first few columns
        ws.append([f'TennisTV.id — {title}'])
        ws.row_dimensions[1].height = 28
        title_cell = ws['A1']
        title_cell.font      = TITLE_FONT
        title_cell.alignment = CENTER
        if headers:
            # Merge across the header columns for the title row
            end_col = openpyxl.utils.get_column_letter(max(len(headers), 1))
            ws.merge_cells(f'A1:{end_col}1')

        # Row 2 — Column headers with green background
        if headers:
            ws.append(headers)
            ws.row_dimensions[2].height = 22
            for col_idx, cell in enumerate(ws[2], start=1):
                cell.fill      = GREEN_FILL
                cell.font      = WHITE_FONT
                cell.alignment = CENTER

        # Rows 3+ — Data (type-safe: keep numbers as numbers, coerce rest to str)
        for row in rows:
            ws.append([safe_cell(cell) for cell in row])

        # Auto-fit column widths
        for col in ws.columns:
            max_len    = 0
            col_letter = col[0].column_letter
            for cell in col:
                try:
                    cell_len = len(str(cell.value)) if cell.value is not None else 0
                    max_len  = max(max_len, cell_len)
                except Exception:
                    pass
            ws.column_dimensions[col_letter].width = min(max(max_len + 4, 10), 50)

        # Freeze rows 1 and 2 (title + header)
        ws.freeze_panes = 'A3'

        output = io.BytesIO()
        wb.save(output)
        output.seek(0)

        return send_file(
            output,
            mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            as_attachment=True,
            download_name='spreadsheet.xlsx',
        )

    except Exception as exc:
        app.logger.error(f'/generate-xlsx error: {exc}')
        return jsonify({'error': str(exc)}), 500


# ══════════════════════════════════════════════════════════════════════════════
#  PDF Generator
#  Brand spec: TennisTV.id header on first page, bold section headings,
#              page numbers in footer on every page
# ══════════════════════════════════════════════════════════════════════════════

class BrandedPDF(FPDF):
    """FPDF subclass that injects a TennisTV.id header and page number footer."""

    def header(self):
        # Only print the brand header on the first page
        if self.page_no() == 1:
            self.set_font('Helvetica', 'B', 10)
            self.set_text_color(29, 106, 58)   # green #1D6A3A
            self.cell(0, 8, 'TennisTV.id', align='L')
            self.set_font('Helvetica', 'I', 9)
            self.set_text_color(150, 150, 150)
            self.cell(0, 8, 'tennistv.id', align='R', ln=True)
            # Divider line
            self.set_draw_color(29, 106, 58)
            self.set_line_width(0.5)
            self.line(10, self.get_y(), 200, self.get_y())
            self.ln(4)

    def footer(self):
        self.set_y(-15)
        self.set_font('Helvetica', 'I', 8)
        self.set_text_color(150, 150, 150)
        self.cell(0, 8, f'Page {self.page_no()} | TennisTV.id', align='C')


@app.route('/generate-pdf', methods=['POST'])
def generate_pdf():
    """
    Expected JSON (matches system prompt schema — sections use "body" field):
    {
      "title": "report title",
      "sections": [
        { "heading": "section heading", "body": "paragraph text" }
      ]
    }
    """
    data = request.get_json(silent=True)
    if not data:
        return jsonify({'error': 'No valid JSON in request body'}), 400

    try:
        pdf = BrandedPDF()
        pdf.set_auto_page_break(auto=True, margin=20)
        pdf.set_margins(left=20, top=25, right=20)
        pdf.add_page()

        DARK  = (26, 26, 46)
        GREEN = (29, 106, 58)
        GRAY  = (100, 100, 100)

        def safe(text):
            """Encode any string to latin-1, replacing unmappable Unicode chars."""
            return str(text).encode('latin-1', errors='replace').decode('latin-1')

        title = safe(data.get('title', 'Report'))

        # ── Document title ──────────────────────────────────────────────────
        pdf.set_font('Helvetica', 'B', 20)
        pdf.set_text_color(*GREEN)
        pdf.multi_cell(0, 11, title, align='C')
        pdf.ln(6)

        # ── Divider ─────────────────────────────────────────────────────────
        pdf.set_draw_color(*GREEN)
        pdf.set_line_width(0.6)
        pdf.line(20, pdf.get_y(), 190, pdf.get_y())
        pdf.ln(8)

        # ── Sections ────────────────────────────────────────────────────────
        for section in data.get('sections', []):
            heading = safe((section.get('heading') or '').strip())
            body    = safe((section.get('body')    or '').strip())

            if heading:
                pdf.set_font('Helvetica', 'B', 13)
                pdf.set_text_color(*GREEN)
                pdf.multi_cell(0, 8, heading)
                pdf.ln(1)

            if body:
                pdf.set_font('Helvetica', '', 11)
                pdf.set_text_color(*DARK)
                pdf.multi_cell(0, 7, body)
                pdf.ln(5)

        pdf_bytes = bytes(pdf.output())
        output = io.BytesIO(pdf_bytes)
        output.seek(0)

        return send_file(
            output,
            mimetype='application/pdf',
            as_attachment=True,
            download_name='report.pdf',
        )

    except Exception as exc:
        app.logger.error(f'/generate-pdf error: {exc}')
        return jsonify({'error': str(exc)}), 500


# ══════════════════════════════════════════════════════════════════════════════
#  Health check
# ══════════════════════════════════════════════════════════════════════════════

@app.route('/', methods=['GET'])
def health():
    return jsonify({'status': 'ok', 'service': 'tennistimes-python-worker'})


if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port, debug=False)
